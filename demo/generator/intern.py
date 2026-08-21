"""Data and DOCX/PDF/PPTX rendering for the "Interne Dienstanweisungen
Meldewesen" library (UPLOAD connector, see docs/features/demo-instance.md).

This is the one library the concept's demo script relies on being visible to
Meldewesen staff (Maria) but invisible to Kfz staff (Thomas) — see
"Berechtigungs-Doppelfrage" in docs/features/demo-instance.md. The rule about
suspected forged identity documents below exists nowhere in the public
Leistungen library on purpose.

Deliberately absent from this whole corpus: anything about a
"Fischereierlaubnis" (fishing licence) — the demo script's "bewusst
unbeantwortbare Frage" needs this topic to have zero coverage anywhere.
"""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO

import docx
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

from zip_utils import normalize_zip_timestamps

SYNTHETIC_NOTICE = (
    "Dieses Dokument ist Teil des synthetischen Demo-Korpus der fiktiven Stadt Rheinfurt "
    "(siehe SOURCE.md im Wurzelverzeichnis dieses Korpus). Es ist nur für die Bibliothek "
    "'Interne Dienstanweisungen Meldewesen' bestimmt und nicht Teil der öffentlichen "
    "Leistungsbeschreibungen."
)


@dataclass
class Dienstanweisung:
    slug: str
    titel: str
    aktenzeichen: str
    abschnitte: list[tuple[str, list[str]]]


@dataclass
class Faq:
    slug: str
    titel: str
    aktenzeichen: str
    fragen: list[tuple[str, str]]


@dataclass
class Folie:
    titel: str
    punkte: list[str]


@dataclass
class Schulung:
    slug: str
    titel: str
    folien: list[Folie]


# --- Dienstanweisungen und Eskalationsregeln (.docx) ------------------------

DIENSTANWEISUNGEN: list[Dienstanweisung] = [
    Dienstanweisung(
        "identitaetszweifel-ausweisantrag",
        "Dienstanweisung: Vorgehen bei Verdacht auf gefälschte oder entwendete "
        "Ausweisdokumente",
        "AZ 32.1-DA-2026-001",
        [
            (
                "Zweck und Geltungsbereich",
                [
                    "Diese Dienstanweisung regelt das Vorgehen der Sachbearbeitung im "
                    "Sachgebiet Meldewesen & Ausweise, wenn beim Beantragen eines "
                    "Personalausweises, Reisepasses oder einer eID-Karte Zweifel an der "
                    "Echtheit vorgelegter Dokumente oder an der Identität der "
                    "antragstellenden Person entstehen.",
                    "Diese Regel ist ausschließlich intern dokumentiert und nicht Bestandteil "
                    "der öffentlichen Leistungsbeschreibungen.",
                ],
            ),
            (
                "Erkennungsmerkmale",
                [
                    "Anzeichen für ein mögliches Sicherheitsrisiko sind unter anderem: "
                    "abweichende Unterschrift gegenüber Altdokument, unplausible "
                    "biografische Angaben, sichtbare Manipulation am Lichtbild oder am "
                    "Sicherheitspapier, sowie Nervosität in Kombination mit fehlenden "
                    "Zusatzdokumenten.",
                ],
            ),
            (
                "Vorgehen",
                [
                    "1. Das Dokument wird nicht zurückgegeben, sondern zunächst einbehalten "
                    "und der Vorgang unterbrochen, ohne die antragstellende Person auf den "
                    "konkreten Verdacht hinzuweisen.",
                    "2. Die Sachgebietsleitung wird unverzüglich telefonisch oder persönlich "
                    "hinzugezogen (Vier-Augen-Prinzip, siehe eigene Dienstanweisung).",
                    "3. Bei erhärtetem Verdacht erfolgt eine Meldung an die Polizeiinspektion "
                    "Rheinfurt über die interne Meldeleitstelle; der Vorgang wird mit "
                    "Aktenzeichen im Vorgangsprotokoll dokumentiert.",
                    "4. Besteht kein hinreichender Verdacht, wird das Dokument zurückgegeben "
                    "und der reguläre Vorgang fortgesetzt.",
                ],
            ),
            (
                "Dokumentation",
                [
                    "Jeder Vorgang nach dieser Dienstanweisung wird im Vorgangsprotokoll des "
                    "Sachgebiets mit Datum, beteiligten Bediensteten und Ergebnis erfasst.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "gebuehrenbefreiung-beduerftigkeit",
        "Dienstanweisung: Gebührenermäßigung und -befreiung wegen Bedürftigkeit",
        "AZ 32.1-DA-2026-002",
        [
            (
                "Rechtsgrundlage",
                [
                    "§ 3 der Verwaltungsgebührensatzung der Stadt Rheinfurt (VGS) erlaubt auf "
                    "schriftlichen Antrag den ganzen oder teilweisen Erlass von "
                    "Verwaltungsgebühren, wenn die Einziehung eine unbillige Härte bedeuten "
                    "würde.",
                ],
            ),
            (
                "Anerkannte Nachweise",
                [
                    "Als Nachweis genügt in der Regel einer der folgenden aktuellen "
                    "Bescheide: Bürgergeld, Grundsicherung im Alter oder bei "
                    "Erwerbsminderung, Wohngeld, oder Leistungen nach dem "
                    "Asylbewerberleistungsgesetz. Der Nachweis darf bei Antragstellung nicht "
                    "älter als drei Monate sein.",
                ],
            ),
            (
                "Verfahren am Schalter",
                [
                    "1. Die antragstellende Person legt den formlosen schriftlichen Antrag "
                    "zusammen mit dem Nachweis vor.",
                    "2. Die Sachbearbeitung prüft die Vollständigkeit des Nachweises und "
                    "leitet den Vorgang zur Entscheidung an die Sachgebietsleitung weiter.",
                    "3. Bis zur Entscheidung wird die Amtshandlung (z. B. Ausstellung des "
                    "Personalausweises) bereits durchgeführt; der Gebührenbescheid ergeht "
                    "gesondert.",
                    "4. Die Entscheidung wird der antragstellenden Person schriftlich "
                    "mitgeteilt und im Vorgang dokumentiert.",
                ],
            ),
            (
                "Hinweis für die Auskunft",
                [
                    "Diese Dienstanweisung konkretisiert § 3 VGS für die tägliche "
                    "Schalterarbeit; die Satzung selbst bleibt die maßgebliche Rechtsgrundlage "
                    "und ist öffentlich einsehbar. Für Kfz-Angelegenheiten gilt § 3 VGS "
                    "entsprechend, eine gesonderte Dienstanweisung für die Kfz-Zulassung "
                    "besteht dazu nicht.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "auskunftssperren-bearbeitung",
        "Dienstanweisung: Bearbeitung von Auskunftssperren nach § 51 BMG",
        "AZ 32.1-DA-2026-003",
        [
            (
                "Zweck",
                [
                    "Diese Dienstanweisung regelt die interne Bearbeitung von Anträgen auf "
                    "Auskunftssperren im Melderegister, insbesondere bei Anhaltspunkten für "
                    "eine Gefährdung der antragstellenden Person.",
                ],
            ),
            (
                "Vorgehen",
                [
                    "Eingehende Anträge werden vorrangig bearbeitet und innerhalb von drei "
                    "Werktagen im Melderegister vermerkt.",
                    "Bei glaubhaft gemachter akuter Gefährdung (z. B. häusliche Gewalt) wird "
                    "die Sperre vorläufig noch am selben Tag eingetragen; die abschließende "
                    "Prüfung der Unterlagen erfolgt nachgelagert.",
                ],
            ),
            (
                "Vertraulichkeit",
                [
                    "Auskünfte über eine bestehende Auskunftssperre werden ausschließlich an "
                    "berechtigte Stellen nach § 51 BMG erteilt; jede Anfrage wird protokolliert.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "vertretungsregelung-meldewesen",
        "Dienstanweisung: Vertretungsregelung im Sachgebiet Meldewesen",
        "AZ 32.1-DA-2026-004",
        [
            (
                "Regelvertretung",
                [
                    "Bei Abwesenheit einer Sachbearbeiterin oder eines Sachbearbeiters "
                    "übernimmt die im Dienstplan ausgewiesene Vertretung innerhalb des "
                    "Sachgebiets die laufenden Vorgänge.",
                ],
            ),
            (
                "Sonderfälle",
                [
                    "Bei gleichzeitiger Abwesenheit von mehr als der Hälfte des Sachgebiets "
                    "wird die Terminvergabe für den betroffenen Zeitraum reduziert; die "
                    "Sachgebietsleitung entscheidet über die Priorisierung dringender Fälle "
                    "(z. B. Reisepassabholung vor gebuchtem Flug).",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "amtshilfe-meldebehoerden",
        "Dienstanweisung: Umgang mit Amtshilfeersuchen anderer Meldebehörden",
        "AZ 32.1-DA-2026-005",
        [
            (
                "Grundsatz",
                [
                    "Ersuchen anderer Meldebehörden werden innerhalb von fünf Werktagen "
                    "beantwortet, sofern keine Auskunftssperre entgegensteht.",
                ],
            ),
            (
                "Formvorgaben",
                [
                    "Ersuchen sind grundsätzlich schriftlich über das Fachverfahren zu "
                    "stellen; telefonische Ersuchen werden nur bei Gefahr im Verzug "
                    "bearbeitet und nachträglich schriftlich bestätigt.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "datenschutz-melderegisterauskunft",
        "Dienstanweisung: Datenschutz bei der Melderegisterauskunft",
        "AZ 32.1-DA-2026-006",
        [
            (
                "Prüfpflicht",
                [
                    "Vor Erteilung einer einfachen Melderegisterauskunft ist zu prüfen, ob "
                    "eine Auskunftssperre, eine Übermittlungssperre oder ein sonstiges "
                    "Übermittlungsverbot vorliegt.",
                ],
            ),
            (
                "Erweiterte Melderegisterauskunft",
                [
                    "Eine erweiterte Auskunft (z. B. an Gerichtsvollzieher) wird nur bei "
                    "Nachweis eines berechtigten Interesses und nach Prüfung durch die "
                    "Sachgebietsleitung erteilt.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "vier-augen-prinzip-ausweisausstellung",
        "Dienstanweisung: Vier-Augen-Prinzip bei der Ausstellung von "
        "Personalausweisen",
        "AZ 32.1-DA-2026-007",
        [
            (
                "Anwendungsfälle",
                [
                    "Das Vier-Augen-Prinzip gilt verbindlich bei Erstbeantragung ohne "
                    "vorhandenes Altdokument, bei Verlustanzeigen mit anschließender "
                    "Neubeantragung am selben Tag, sowie bei jedem Verdachtsfall nach der "
                    "Dienstanweisung zu gefälschten Ausweisdokumenten.",
                ],
            ),
            (
                "Durchführung",
                [
                    "Eine zweite Sachbearbeiterin oder ein zweiter Sachbearbeiter prüft "
                    "unabhängig die vorgelegten Unterlagen und bestätigt die Freigabe mit "
                    "Kurzzeichen im Fachverfahren, bevor der Antrag weiterverarbeitet wird.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "verlust-dienstlicher-zugangsdaten",
        "Dienstanweisung: Verfahren bei Verlust dienstlicher Zugangsdaten",
        "AZ 32.1-DA-2026-008",
        [
            (
                "Sofortmaßnahmen",
                [
                    "Der Verlust von dienstlichen Zugangsdaten oder eines Dienstausweises ist "
                    "unverzüglich der Sachgebietsleitung und der IT-Leitstelle zu melden; "
                    "betroffene Zugänge werden umgehend gesperrt.",
                ],
            ),
            (
                "Dokumentation",
                [
                    "Der Vorfall wird im internen Meldeportal dokumentiert; die IT-Leitstelle "
                    "prüft innerhalb von 24 Stunden auf Auffälligkeiten in den betroffenen "
                    "Konten.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "terminvergabe-wartezeitmanagement",
        "Dienstanweisung: Terminvergabe und Wartezeitmanagement im Bürgerbüro",
        "AZ 32.1-DA-2026-009",
        [
            (
                "Kontingente",
                [
                    "Je Sachgebiet werden täglich feste Terminkontingente sowie ein kleines "
                    "Kontingent für dringende Spontanfälle ohne Termin vorgehalten.",
                ],
            ),
            (
                "Priorisierung",
                [
                    "Bei hohem Andrang entscheidet die diensthabende Teamleitung über eine "
                    "vorübergehende Umverteilung von Personal zwischen den "
                    "Empfangsbereichen Meldewesen und Kfz-Zulassung.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "umgang-aggressives-verhalten",
        "Dienstanweisung: Umgang mit aggressivem oder bedrohlichem Verhalten am "
        "Schalter",
        "AZ 32.1-DA-2026-010",
        [
            (
                "Deeskalation",
                [
                    "Bei erkennbarer Eskalation wird zunächst ruhig und sachlich deeskaliert; "
                    "Hilfe durch Kolleginnen und Kollegen ist frühzeitig, aber unauffällig "
                    "anzufordern.",
                ],
            ),
            (
                "Eskalationsstufen",
                [
                    "Bleibt die Deeskalation erfolglos, gelten die Eskalationsregeln Stufe 1 "
                    "bis 3 bei Beschwerden im Bürgerbüro; bei unmittelbarer Gefahr wird der "
                    "Empfang über den stillen Alarm benachrichtigt und die Polizei gerufen.",
                ],
            ),
        ],
    ),
]

ESKALATIONSREGELN: list[Dienstanweisung] = [
    Dienstanweisung(
        "eskalation-beschwerden-buergerbuero",
        "Eskalationsregeln Stufe 1 bis 3 bei Beschwerden im Bürgerbüro",
        "AZ 32.1-ESK-2026-001",
        [
            (
                "Stufe 1: Vor-Ort-Klärung",
                [
                    "Die Sachbearbeitung versucht, die Beschwerde direkt am Schalter zu "
                    "klären. Gelingt dies nicht innerhalb weniger Minuten, wird an Stufe 2 "
                    "übergeben.",
                ],
            ),
            (
                "Stufe 2: Teamleitung",
                [
                    "Die diensthabende Teamleitung übernimmt das Gespräch. Bleibt die "
                    "Beschwerde ungelöst oder betrifft sie einen grundsätzlichen "
                    "Verfahrensfehler, wird an Stufe 3 übergeben.",
                ],
            ),
            (
                "Stufe 3: Sachgebietsleitung",
                [
                    "Die Sachgebietsleitung entscheidet abschließend und dokumentiert die "
                    "Beschwerde im Beschwerdemanagement-System.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "eskalation-it-ausfall",
        "Eskalationsregeln bei IT-Ausfall der Fachanwendung MeldeSoft",
        "AZ 32.1-ESK-2026-002",
        [
            (
                "Sofortmaßnahmen",
                [
                    "Bei Ausfall der Fachanwendung MeldeSoft wechselt das Sachgebiet auf den "
                    "papiergestützten Notbetrieb; laufende Vorgänge werden mit Zeitstempel auf "
                    "Papier fortgeführt und nach Wiederherstellung nacherfasst.",
                ],
            ),
            (
                "Meldewege",
                [
                    "Der Ausfall wird unverzüglich der IT-Leitstelle gemeldet; dauert die "
                    "Störung länger als 30 Minuten, informiert die Sachgebietsleitung die "
                    "wartenden Bürgerinnen und Bürger über die voraussichtliche Dauer.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "eskalation-kindeswohlgefaehrdung",
        "Eskalationsregeln bei Verdacht auf Kindeswohlgefährdung während der "
        "Vorsprache",
        "AZ 32.1-ESK-2026-003",
        [
            (
                "Erkennen",
                [
                    "Ergeben sich während einer Vorsprache konkrete Anhaltspunkte für eine "
                    "Kindeswohlgefährdung, wird das Gespräch besonnen fortgeführt, ohne die "
                    "Situation zu verschärfen.",
                ],
            ),
            (
                "Meldung",
                [
                    "Die Sachgebietsleitung wird unverzüglich informiert und meldet den "
                    "Sachverhalt an die Fachberatung Kinderschutz gemäß den städtischen "
                    "Kinderschutzverfahren.",
                ],
            ),
        ],
    ),
    Dienstanweisung(
        "eskalation-medizinischer-notfall",
        "Eskalationsregeln bei medizinischen Notfällen im Bürgerbüro",
        "AZ 32.1-ESK-2026-004",
        [
            (
                "Erstversorgung",
                [
                    "Ersthelfende Kolleginnen und Kollegen leisten Erste Hilfe; parallel wird "
                    "der Rettungsdienst über die zentrale Notrufnummer alarmiert.",
                ],
            ),
            (
                "Nachbereitung",
                [
                    "Der Vorfall wird im Meldeportal für Arbeitssicherheit dokumentiert; bei "
                    "Bedarf wird betroffenen Mitarbeitenden ein Gespräch mit dem "
                    "betriebsärztlichen Dienst angeboten.",
                ],
            ),
        ],
    ),
]


def render_dienstanweisung_docx(da: Dienstanweisung) -> bytes:
    document = docx.Document()
    style = document.styles["Normal"]
    style.font.size = Pt(11)

    document.add_heading(da.titel, level=1)
    meta = document.add_paragraph()
    meta.add_run(f"Aktenzeichen: {da.aktenzeichen}").italic = True

    for ueberschrift, absaetze in da.abschnitte:
        document.add_heading(ueberschrift, level=2)
        for absatz in absaetze:
            document.add_paragraph(absatz)

    document.add_paragraph()
    notice = document.add_paragraph()
    notice_run = notice.add_run(SYNTHETIC_NOTICE)
    notice_run.italic = True
    notice_run.font.size = Pt(9)

    buffer = BytesIO()
    document.save(buffer)
    return normalize_zip_timestamps(buffer.getvalue())


# --- Interne FAQ (.pdf) ------------------------------------------------------

STYLES = getSampleStyleSheet()
FAQ_TITLE_STYLE = ParagraphStyle("FaqTitle", parent=STYLES["Title"], fontSize=15, spaceAfter=12)
FAQ_QUESTION_STYLE = ParagraphStyle(
    "FaqQuestion", parent=STYLES["Heading3"], spaceBefore=10, spaceAfter=2
)
FAQ_ANSWER_STYLE = ParagraphStyle("FaqAnswer", parent=STYLES["BodyText"], spaceAfter=6)
FAQ_FOOTER_STYLE = ParagraphStyle(
    "FaqFooter", parent=STYLES["BodyText"], fontSize=8, textColor=colors.grey
)

FAQS: list[Faq] = [
    Faq(
        "faq-ausweisbeantragung",
        "Interne FAQ: Häufige Rückfragen zur Ausweisbeantragung",
        "AZ 32.1-FAQ-2026-001",
        [
            (
                "Darf ein Ausweis auch ohne Termin ausgestellt werden?",
                "Nur im Rahmen des Spontan-Kontingents nach der Dienstanweisung zur "
                "Terminvergabe, und nur bei glaubhaft gemachter Dringlichkeit (z. B. "
                "gebuchte Reise innerhalb von drei Werktagen).",
            ),
            (
                "Was tun, wenn das biometrische Passbild Zweifel weckt?",
                "Bei sichtbaren Bearbeitungsspuren wird ein neues Passbild verlangt; "
                "bestehen darüber hinaus Zweifel an der Identität, greift die Dienstanweisung "
                "zu gefälschten Ausweisdokumenten.",
            ),
            (
                "Kann die Gebühr gestundet werden?",
                "Eine Stundung ist nicht vorgesehen; bei Bedürftigkeit gilt stattdessen § 3 "
                "VGS in Verbindung mit der entsprechenden Dienstanweisung.",
            ),
        ],
    ),
    Faq(
        "faq-ummeldung",
        "Interne FAQ: Häufige Rückfragen zur Ummeldung",
        "AZ 32.1-FAQ-2026-002",
        [
            (
                "Was, wenn die Wohnungsgeberbestätigung fehlt?",
                "Der Vorgang wird vorläufig mit Fristsetzung von zwei Wochen zur "
                "Nachreichung angelegt; ohne Nachreichung erfolgt ein Bußgeldhinweis nach "
                "dem Bundesmeldegesetz.",
            ),
            (
                "Muss die ganze Familie gemeinsam vorsprechen?",
                "Nein, eine bevollmächtigte Person kann für Familienangehörige mitvorsprechen, "
                "sofern für jede Person eine unterschriebene Vollmacht vorliegt.",
            ),
        ],
    ),
    Faq(
        "faq-kfz-schalterfragen",
        "Interne FAQ: Häufige Rückfragen zur Kfz-Zulassung am gemeinsamen Empfang",
        "AZ 32.1-FAQ-2026-003",
        [
            (
                "Wohin verweise ich Rückfragen zum Wunschkennzeichen?",
                "Die Reservierung läuft über das Online-Portal der Kfz-Zulassungsstelle; am "
                "Meldewesen-Schalter wird nur auf das Portal und die Öffnungszeiten der "
                "Kfz-Zulassungsstelle verwiesen.",
            ),
            (
                "Darf ich als Meldewesen-Sachbearbeitung Kfz-Vorgänge bearbeiten?",
                "Nein, die Bearbeitung von Kfz-Vorgängen ist dem Sachgebiet Kfz-Zulassung "
                "vorbehalten; bei Personalengpässen entscheidet die Teamleitung über eine "
                "vorübergehende Umverteilung.",
            ),
        ],
    ),
    Faq(
        "faq-fundsachen-empfang",
        "Interne FAQ: Umgang mit Fundsachen am Empfang",
        "AZ 32.1-FAQ-2026-004",
        [
            (
                "Wohin mit Fundsachen, die am Meldewesen-Schalter abgegeben werden?",
                "Fundsachen werden mit Fundzettel versehen und noch am selben Tag an das "
                "städtische Fundbüro weitergeleitet.",
            ),
            (
                "Wie lange verwahrt das Fundbüro Gegenstände?",
                "Gemäß Fundgebührensatzung in der Regel sechs Monate; danach kann die Stadt "
                "Rheinfurt die Sache verwerten oder versteigern.",
            ),
        ],
    ),
    Faq(
        "faq-presseanfragen-schalter",
        "Interne FAQ: Verhalten bei Presseanfragen am Schalter",
        "AZ 32.1-FAQ-2026-005",
        [
            (
                "Darf ich Presseanfragen direkt beantworten?",
                "Nein, Presseanfragen werden grundsätzlich an das Presseamt der Stadt "
                "Rheinfurt weitergeleitet (presse@stadt-rheinfurt.example).",
            ),
            (
                "Was, wenn ein Kamerateam ohne Anmeldung erscheint?",
                "Die Teamleitung wird informiert; Aufnahmen im laufenden Publikumsverkehr "
                "sind ohne vorherige Drehgenehmigung nicht gestattet.",
            ),
        ],
    ),
    Faq(
        "faq-standesamt-zusammenarbeit",
        "Interne FAQ: Zusammenarbeit mit dem Standesamt",
        "AZ 32.1-FAQ-2026-006",
        [
            (
                "Wer informiert das Standesamt über eine Eheschließung?",
                "Das Standesamt informiert die Meldebehörde automatisch über das "
                "Fachverfahren; eine gesonderte Meldung durch das Meldewesen ist nicht "
                "erforderlich.",
            ),
            (
                "Wohin mit Rückfragen zu Ehefähigkeitszeugnissen?",
                "Diese werden ausschließlich vom Standesamt bearbeitet und dorthin "
                "weiterverwiesen.",
            ),
        ],
    ),
]


def render_faq_pdf(faq: Faq) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=2.2 * cm,
        rightMargin=2.2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        title=faq.titel,
        author="Bürgerbüro Rheinfurt (synthetisch)",
    )
    story = [
        Paragraph(faq.titel, FAQ_TITLE_STYLE),
        Paragraph(f"Aktenzeichen (Muster): {faq.aktenzeichen}", FAQ_FOOTER_STYLE),
        Spacer(1, 0.4 * cm),
    ]
    for frage, antwort in faq.fragen:
        story.append(Paragraph(frage, FAQ_QUESTION_STYLE))
        story.append(Paragraph(antwort, FAQ_ANSWER_STYLE))
    story.append(Spacer(1, 0.6 * cm))
    story.append(Paragraph(SYNTHETIC_NOTICE, FAQ_FOOTER_STYLE))
    doc.build(story)
    return buffer.getvalue()


# --- Schulungsfolien (.pptx) -------------------------------------------------

SCHULUNGEN: list[Schulung] = [
    Schulung(
        "onboarding-buergerbuero",
        "Onboarding neue Mitarbeitende Bürgerbüro",
        [
            Folie("Onboarding neue Mitarbeitende Bürgerbüro", ["Sachgebiet Meldewesen & Ausweise"]),
            Folie(
                "Struktur des Bürgerbüros",
                [
                    "Sachgebiete: Meldewesen & Ausweise, Kfz-Zulassung",
                    "Amtsleitung Bürgerbüro",
                    "Gemeinsamer Empfang beider Sachgebiete",
                ],
            ),
            Folie(
                "Wichtige Dienstanweisungen",
                [
                    "Vier-Augen-Prinzip bei Ausweisausstellung",
                    "Vorgehen bei Verdacht auf gefälschte Dokumente",
                    "Gebührenermäßigung wegen Bedürftigkeit",
                ],
            ),
            Folie(
                "Erste Woche",
                [
                    "Einweisung in MeldeSoft",
                    "Begleitete Schalterdienste",
                    "Kennenlernen der Eskalationsregeln",
                ],
            ),
        ],
    ),
    Schulung(
        "betrugserkennung-ausweisdokumente",
        "Betrugserkennung bei Ausweisdokumenten (Basisschulung)",
        [
            Folie("Betrugserkennung bei Ausweisdokumenten", ["Basisschulung Sachgebiet Meldewesen"]),
            Folie(
                "Sicherheitsmerkmale",
                [
                    "Hologramm und Kippeffekt",
                    "Lasergravur der biografischen Daten",
                    "Mikroschrift im Untergrund",
                ],
            ),
            Folie(
                "Typische Auffälligkeiten",
                [
                    "Abweichende Unterschrift zum Altdokument",
                    "Manipulationsspuren am Lichtbild",
                    "Unplausible biografische Angaben",
                ],
            ),
            Folie(
                "Vorgehen bei Verdacht",
                [
                    "Dokument einbehalten, Vorgang unterbrechen",
                    "Sachgebietsleitung hinzuziehen",
                    "Bei erhärtetem Verdacht: Meldung an die Polizeiinspektion",
                ],
            ),
        ],
    ),
    Schulung(
        "datenschutzgrundlagen-meldewesen",
        "Datenschutzgrundlagen für Meldewesen",
        [
            Folie("Datenschutzgrundlagen für Meldewesen", []),
            Folie(
                "Rechtsgrundlagen",
                ["Bundesmeldegesetz (BMG)", "Datenschutz-Grundverordnung (DSGVO)"],
            ),
            Folie(
                "Auskunftssperren",
                [
                    "§ 51 BMG: einfache und erweiterte Sperre",
                    "Vorrangige Bearbeitung bei Gefährdungslage",
                ],
            ),
            Folie(
                "Melderegisterauskünfte",
                [
                    "Immer Sperrvermerke prüfen",
                    "Erweiterte Auskunft nur mit berechtigtem Interesse",
                ],
            ),
        ],
    ),
    Schulung(
        "deeskalationstraining-schalter",
        "Deeskalationstraining am Schalter",
        [
            Folie("Deeskalationstraining am Schalter", []),
            Folie(
                "Warnsignale erkennen",
                ["Erhöhte Lautstärke", "Angespannte Körpersprache", "Wiederholte Unterbrechungen"],
            ),
            Folie(
                "Deeskalationstechniken",
                ["Ruhiger Tonfall", "Aktives Zuhören", "Klare, kurze Sätze"],
            ),
            Folie(
                "Wann eskalieren?",
                [
                    "Eskalationsregeln Stufe 1 bis 3 anwenden",
                    "Bei Gefahr: stiller Alarm und Polizei",
                ],
            ),
        ],
    ),
    Schulung(
        "barrierefreie-kommunikation",
        "Barrierefreie Kommunikation im Bürgerbüro",
        [
            Folie("Barrierefreie Kommunikation im Bürgerbüro", []),
            Folie(
                "Leichte Sprache",
                ["Kurze Sätze", "Ein Gedanke pro Satz", "Fremdwörter vermeiden"],
            ),
            Folie(
                "Hör- und Sehbeeinträchtigungen",
                [
                    "Braille-Aufkleber für den Personalausweis anbieten",
                    "Bei Bedarf Schriftdolmetschung organisieren",
                ],
            ),
            Folie(
                "Mehrsprachigkeit",
                ["Dolmetschliste des Bürgerbüros nutzen", "Wichtige Formulare in mehreren Sprachen"],
            ),
        ],
    ),
    Schulung(
        "fachanwendung-meldesoft",
        "Umgang mit der Fachanwendung MeldeSoft",
        [
            Folie("Umgang mit der Fachanwendung MeldeSoft", []),
            Folie(
                "Grundfunktionen",
                ["Vorgang anlegen", "Dokumente scannen und verknüpfen", "Gebühren erfassen"],
            ),
            Folie(
                "Bei Systemausfall",
                [
                    "Wechsel in den papiergestützten Notbetrieb",
                    "IT-Leitstelle informieren",
                    "Nacherfassung nach Wiederherstellung",
                ],
            ),
            Folie(
                "Datenqualität",
                ["Adressfelder immer vollständig pflegen", "Dubletten vor Speichern prüfen"],
            ),
        ],
    ),
]


def render_schulung_pptx(schulung: Schulung) -> bytes:
    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]

    first = schulung.folien[0]
    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = first.titel
    if first.punkte:
        title_slide.placeholders[1].text = first.punkte[0]

    for folie in schulung.folien[1:]:
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = folie.titel
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, punkt in enumerate(folie.punkte):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = punkt
            paragraph.font.size = PptxPt(20)

    # Closing slide with the synthetic-data notice, so every generated
    # deck carries it just like the docx/pdf documents do.
    closing = presentation.slides.add_slide(bullet_layout)
    closing.shapes.title.text = "Hinweis"
    body = closing.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = SYNTHETIC_NOTICE
    body.paragraphs[0].font.size = PptxPt(14)

    buffer = BytesIO()
    presentation.save(buffer)
    return normalize_zip_timestamps(buffer.getvalue())
